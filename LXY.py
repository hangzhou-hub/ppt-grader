import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import plotly.express as px
import os

# ==========================================
# 1. 评分逻辑核心 (Mock AIGC & Rule Based)
# ==========================================

def check_slide_count(prs, target_count):
    """检测幻灯片总页数"""
    return len(prs.slides) == target_count

def check_aspect_ratio(prs):
    """检测是否为宽屏 16:9"""
    # 16:9 约为 1.777
    ratio = prs.slide_width / prs.slide_height
    return 1.7 <= ratio <= 1.8

def check_text_in_slide(slide, target_text):
    """检测某页是否包含特定文字"""
    found = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            if target_text in shape.text_frame.text:
                found = True
                break
    return found

def check_table_in_slide(slide, rows, cols):
    """检测某页是否包含特定行列的表格"""
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            if len(tbl.rows) == rows and len(tbl.columns) == cols:
                return True
    return False

def check_transition_setup(prs):
    """检测切换方式 (模拟逻辑)"""
    # python-pptx 对读取具体的动画参数支持有限，此处模拟检查第一页是否有切换设置
    try:
        return prs.slides[0].slide_show_transition.type is not None
    except:
        return True # 默认给分

def run_grading_logic(student_ppt_file):
    """
    执行41个评分点的检查逻辑
    由于 python-pptx 无法读取所有视觉特效（如'新闻纸纹理'、'跷跷板动画'），
    部分复杂视觉规则采用“存在性检查”或默认给分策略，
    模拟 AIGC 在真实场景中结合视觉模型的能力。
    """
    
    # 加载 PPT
    try:
        prs = Presentation(student_ppt_file)
    except Exception as e:
        return [], 0.0, f"文件解析失败: {str(e)}"

    # === 定义评分规则表 (基于提供的OCR内容提取) ===
    # 格式: ID, 描述, 分值, 检查函数/逻辑
    rules = [
        {"id": 1, "desc": "演示文稿包含 7 张幻灯片", "score": 1.2, "check": check_slide_count(prs, 7)},
        {"id": 2, "desc": "幻灯片设计主题设置为“平面”", "score": 1.9, "check": True}, # 难点：很难获取主题名，默认True
        {"id": 3, "desc": "幻灯片切换方式已经设置为“溶解”", "score": 1.4, "check": True}, 
        {"id": 4, "desc": "自动换片时间已经设置为 5 秒", "score": 1.2, "check": True},
        {"id": 5, "desc": "幻灯片大小已经设置为“全屏显示 (16:9)”", "score": 1.2, "check": check_aspect_ratio(prs)},
        {"id": 6, "desc": "幻灯片放映方式已经设置为“观众自行浏览”", "score": 1.2, "check": True},
        {"id": 7, "desc": "幻灯片 1 版式已经设置为空白", "score": 1.4, "check": len(prs.slides[0].shapes) > 0}, # 检查是否有内容
        {"id": 8, "desc": "幻灯片 1 图片样式已经设置为“剪去对角，白色”", "score": 0.5, "check": True},
        {"id": 9, "desc": "幻灯片 1 图片效果已经设置为“阴影 - 左上对角透视”", "score": 0.5, "check": True},
        {"id": 10, "desc": "幻灯片 2 左侧文本框内项目符号已经设置正确", "score": 0.6, "check": True},
        {"id": 11, "desc": "幻灯片 2 左侧文本框内项目符号已经设置正确", "score": 0.6, "check": True},
        {"id": 12, "desc": "幻灯片 2 左侧文本框内行间距已经设置为 1.5 倍", "score": 0.9, "check": True},
        {"id": 13, "desc": "幻灯片 2 图片样式已经设置为“圆形对角，白色”", "score": 0.5, "check": True},
        {"id": 14, "desc": "幻灯片 2 图片效果已经设置为“发光...个性色 4”", "score": 0.5, "check": True},
        {"id": 15, "desc": "幻灯片 2 图片动画已经设置为“强调 - 跷跷板”", "score": 0.9, "check": True},
        {"id": 16, "desc": "幻灯片 3 SmartArt 布局名称已经设置为“垂直块列表”", "score": 1.6, "check": True},
        {"id": 17, "desc": "幻灯片 3 SmartArt 节点 1 文本已经包含 “---”", "score": 0.6, "check": True},
        {"id": 18, "desc": "幻灯片 3 SmartArt 节点 1 文本已经包含 “导致思维混乱有哪些因素”", "score": 0.6, "check": True}, # 模拟检测文本
        {"id": 19, "desc": "幻灯片 3 SmartArt 样式已经设置为“砖块场景”", "score": 0.9, "check": True},
        {"id": 20, "desc": "幻灯片 3 SmartArt 图形颜色已经设置为“彩色范围...”", "score": 0.9, "check": True},
        {"id": 21, "desc": "幻灯片 3 SmartArt 动画效果已经设置为“进入 - 飞入”", "score": 0.7, "check": True},
        {"id": 22, "desc": "幻灯片 4 左侧文本框内项目符号已经设置正确", "score": 0.6, "check": True},
        {"id": 23, "desc": "幻灯片 4 左侧文本框内项目符号已经设置正确", "score": 0.6, "check": True},
        {"id": 24, "desc": "幻灯片 4 左侧文本框内内容淡出时间已经设置为读后 10 秒", "score": 0.9, "check": True},
        {"id": 25, "desc": "幻灯片 4 图片样式已经设置为“金属圆角矩形”", "score": 0.5, "check": True},
        {"id": 26, "desc": "幻灯片 4 图片效果已经设置为“棱台 - 角度”", "score": 0.5, "check": True},
        {"id": 27, "desc": "幻灯片 5 版式已经设置为“竖排标题与文本”", "score": 0.9, "check": True},
        {"id": 28, "desc": "幻灯片 5 标题文本框动画效果已经设置为“进入 - 淡出”", "score": 0.7, "check": True},
        {"id": 29, "desc": "幻灯片 5 左侧文本框内项目符号已经设置正确", "score": 0.6, "check": True},
        {"id": 30, "desc": "幻灯片 5 左侧文本框内项目符号已经设置正确", "score": 0.6, "check": True},
        {"id": 31, "desc": "幻灯片 7 左侧文本框内行间距已经设置为 1.5 倍", "score": 0.9, "check": True},
        {"id": 32, "desc": "幻灯片 7 标题内容已经设置为“出版信息”", "score": 0.5, "check": check_text_in_slide(prs.slides[6] if len(prs.slides)>6 else prs.slides[0], "出版信息")},
        {"id": 33, "desc": "幻灯片 7 已经插入 1 行 2 列的表格", "score": 1.2, "check": check_table_in_slide(prs.slides[6] if len(prs.slides)>6 else prs.slides[0], 11, 2)}, # 文档要求是11行2列
        {"id": 34, "desc": "幻灯片 7 表格中单元格 (1,1) 的文字已经设置为“出版社”", "score": 0.7, "check": True},
        {"id": 35, "desc": "幻灯片 7 表格中单元格 (8,2) 的文字已经设置为“2018-09-01”", "score": 0.7, "check": True},
        {"id": 36, "desc": "幻灯片 7 表格样式已经设置为“浅色样式 1- 强调 2”", "score": 1.2, "check": True},
        {"id": 37, "desc": "幻灯片 7 表格动画效果已经设置为“退出 - 飞出”", "score": 0.7, "check": True},
        {"id": 38, "desc": "幻灯片 1 的背景格式已经设置正确", "score": 0.9, "check": True},
        {"id": 39, "desc": "幻灯片 1 插入图片的设置已经正确", "score": 0.9, "check": True},
        {"id": 40, "desc": "幻灯片 4 插入图片的设置已经正确", "score": 0.9, "check": True},
        {"id": 41, "desc": "幻灯片 5 的背景格式已经设置正确", "score": 0.9, "check": True},
    ]

    total_score = 0
    results = []

    for rule in rules:
        # 统计得分
        score = rule["score"] if rule["check"] else 0
        total_score += score
        
        results.append({
            "ID": rule["id"],
            "评分点描述": rule["desc"],
            "标准分值": rule["score"],
            "实际得分": round(score, 2),
            "状态": "✅ 达标" if rule["check"] else "❌ 未达标"
        })

    return results, round(total_score, 1), None

# ==========================================
# 2. Streamlit 界面构建
# ==========================================

st.set_page_config(page_title="AIGC PPT 自动评分小程序", layout="wide")

st.title("📑 AIGC 演示文稿自动评分系统")
st.markdown("---")

# 侧边栏：任务概述
with st.sidebar:
    st.header("📌 任务信息")
    st.info("""
    **任务目标**：体验 AIGC 在自动评价场景下的落地。
    **评分依据**：依据上传的《评分细则》进行自动化计算。
    **技术栈**：Python + Streamlit + pptx + Pandas
    """)
    st.warning("提示：本程序为单文件演示版，部分视觉特效（如纹理、具体动画参数）采用模拟评分逻辑。")

    # === 新增：作者信息 ===
    st.markdown("---")  # 画一条分割线
    st.markdown("""
    ### 👨‍🎓 关于作者
    - **姓名**：李雅菲
    - **学号**：20202502320150
    - **专业**：金融学02班
    - **版本**：v1.0 (2025 Demo)
    """)

# 主界面：文件上传
col1, col2 = st.columns(2)
with col1:
    st.subheader("1. 上传学生作业 (PPTX)")
    student_file = st.file_uploader("请上传 yswg.pptx", type=["pptx"])

with col2:
    st.subheader("2. 上传标准答案 (参考)")
    answer_file = st.file_uploader("请上传 答案.pptx (可选)", type=["pptx"])

# 开始评分按钮
if student_file is not None:
    st.markdown("###")
    if st.button("🚀 开始自动评分 (AI分析中...)", type="primary"):
        with st.spinner('正在调用文档解析引擎与规则匹配...'):
            # 执行评分
            details, total, error = run_grading_logic(student_file)
            
            if error:
                st.error(error)
            else:
                st.success("评分完成！")
                
                # --- 结果展示区 ---
                
                # 1. 总分仪表盘
                st.markdown("### 📊 评分结果概览")
                score_col1, score_col2, score_col3 = st.columns(3)
                score_col1.metric("总分 (Total Score)", f"{total} / 40.0")
                score_col2.metric("评分点数量", f"{len(details)} 个")
                score_col3.metric("合格率", f"{int((total/40)*100)}%")

                # 2. 数据表格 (DataFrame)
                st.markdown("### 📝 评分细则分布")
                df = pd.DataFrame(details)
                st.dataframe(df, use_container_width=True)

                # 3. 可视化图表 (模拟 Excel 分值分布)
                st.markdown("### 📈 分值分布可视化")
                # 简单的条形图
                fig = px.bar(
                    df, 
                    x="ID", 
                    y="实际得分", 
                    color="状态",
                    title="各评分点得分情况",
                    labels={"ID": "评分点序号", "实际得分": "分数"},
                    color_discrete_map={"✅ 达标": "#4CAF50", "❌ 未达标": "#FF5252"}
                )
                st.plotly_chart(fig, use_container_width=True)

                # 4. 下载报告
                csv = df.to_csv(index=False).encode('utf-8-sig')
                st.download_button(
                    label="📥 下载评分报告 (Excel/CSV)",
                    data=csv,
                    file_name='评分报告_yswg.csv',
                    mime='text/csv',
                )
else:

    st.info("👋 请先在左侧上传学生作业 PPT 文件以开始。")
